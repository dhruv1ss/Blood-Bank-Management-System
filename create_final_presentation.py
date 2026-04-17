from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_RED = RGBColor(220, 38, 38)
DARK_BG = RGBColor(20, 20, 30)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_GREEN = RGBColor(34, 197, 94)
ACCENT_PURPLE = RGBColor(139, 92, 246)
ACCENT_ORANGE = RGBColor(249, 115, 22)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add decorative shapes
    shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = PRIMARY_RED
    shape1.line.color.rgb = PRIMARY_RED
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_paragraph_slide(prs, title, paragraph_text, accent_color=PRIMARY_RED):
    """Add a slide with paragraph text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = accent_color
    title_shape.line.color.rgb = accent_color
    
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    
    # Add paragraph content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = paragraph_text
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BG
    p.line_spacing = 1.5
    
    return slide

def add_image_slide(prs, title, image_path, description, accent_color=ACCENT_BLUE):
    """Add a slide with image and description"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = accent_color
    title_shape.line.color.rgb = accent_color
    
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    
    # Add image if it exists
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.1), width=Inches(9), height=Inches(4.5))
        except:
            pass
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8.6), Inches(1.4))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BG
    p.line_spacing = 1.3
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "LifeFlow", "🩸 Blood Donation Management System 🩸")

# Slide 2: Executive Summary
add_paragraph_slide(prs, "Executive Summary", 
"""LifeFlow is a comprehensive digital platform designed to revolutionize blood donation management in the modern era. Our system connects donors, organizations, and patients in a seamless ecosystem that solves critical blood shortage problems through innovative technology and gamification. The platform provides real-time matching, tracking, and incentive systems that encourage regular donations and improve emergency response times. LifeFlow is a scalable solution suitable for hospitals, blood banks, NGOs, and healthcare organizations worldwide. By leveraging modern web technologies and user-centric design, we aim to save lives and build a community of dedicated blood donors.""", PRIMARY_RED)

# Slide 3: Problem Statement
add_paragraph_slide(prs, "The Problem We Solve",
"""Blood shortage is a critical global health crisis affecting millions of patients every year. Hospitals and medical centers struggle to maintain adequate blood supplies, especially during emergencies. The current system lacks a centralized platform for blood donation management, making it difficult to find willing donors when needed. There is no transparency in donation tracking and verification, and limited awareness about donation camps and opportunities. Additionally, there is no incentive system to encourage regular donors, resulting in inconsistent donation rates. LifeFlow addresses all these challenges by providing a comprehensive, user-friendly platform that connects all stakeholders in the blood donation ecosystem.""", ACCENT_ORANGE)

# Slide 4: Our Solution
add_paragraph_slide(prs, "Our Solution",
"""LifeFlow provides a centralized blood donation management platform that revolutionizes how blood donations are organized and tracked. The system enables easy donor registration with blood group and location tracking, allowing organizations to quickly identify potential donors. Our real-time blood request matching system connects patients in need with available donors instantly. The location-based donation camp discovery feature with interactive maps helps donors find nearby camps easily. We incorporate gamification with badges, points, and leaderboards to motivate regular donations. Finally, our comprehensive admin panel allows system administrators to approve requests, manage users, and oversee the entire platform. Together, these features create an efficient, transparent, and engaging blood donation ecosystem.""", ACCENT_GREEN)

# Slide 5: Key Features
add_paragraph_slide(prs, "Key Features Overview",
"""LifeFlow includes several powerful features designed to streamline blood donation management. Secure Authentication ensures user data protection through email/password login with JWT tokens. The Donation Camps feature provides an interactive map showing camp locations, dates, times, and available slots. The Blood Requests system allows patients and organizations to post emergency blood requests that are matched with available donors. The Heroes Leaderboard gamifies the donation experience with badges, points, and rankings to recognize and reward top donors. The Admin Dashboard provides complete system management, analytics, and request approval capabilities. Finally, the Organization Panel allows blood banks and NGOs to create and manage donation camps, track registrations, and view statistics. Each feature is designed with user experience and efficiency in mind.""", ACCENT_BLUE)

# Slide 6: Home Page
add_image_slide(prs, "Home Page - Welcome & Engagement",
    "screenshots/01_home_page.png",
    """The Home Page serves as the entry point to LifeFlow, featuring a compelling hero section with clear call-to-action buttons. It displays impact statistics showing the number of lives saved, active donors, and successful donations. The page includes a featured heroes section highlighting top donors and their achievements. A 'How It Works' section explains the donation process step-by-step. The design is fully responsive and optimized for all devices, with smooth animations and modern UI elements that engage visitors and encourage them to participate in the blood donation movement.""", ACCENT_BLUE)

# Slide 7: Login Page
add_image_slide(prs, "Login Page - Secure Authentication",
    "screenshots/02_login_page.png",
    """The Login Page provides secure access to the LifeFlow platform with email and password authentication. The page features a modern glassmorphism design with a semi-transparent background and blur effects. Users can enter their email and password, with real-time validation and error messages. The page includes a 'Remember Me' checkbox for convenient access and quick links to the registration page for new users. The responsive design ensures smooth login experience on mobile and desktop devices. Behind the scenes, the system uses JWT tokens for secure session management and Bcrypt for password hashing.""", ACCENT_BLUE)

# Slide 8: Registration Page
add_image_slide(prs, "Registration Page - User Onboarding",
    "screenshots/03_register_page.png",
    """The Registration Page guides new users through the onboarding process with a comprehensive form. Users provide their full name, email, password, blood group, age, and location information. The form includes role selection allowing users to register as Donors or Organizations. For organizations, additional fields collect organization name, contact phone, and address. The page features real-time form validation with helpful error messages. Password requirements are clearly displayed to ensure security. The responsive design makes registration easy on all devices. Upon successful registration, users are automatically logged in and redirected to their dashboard.""", ACCENT_BLUE)

# Slide 9: Donation Camps Page
add_image_slide(prs, "Donation Camps - Location Discovery",
    "screenshots/04_donation_camps.png",
    """The Donation Camps Page features an interactive map showing all available blood donation camps in real-time. Users can zoom, pan, and click on camp markers to view detailed information including camp name, date, time, address, and available slots. The page includes filtering options to search camps by location, date, and blood types needed. Each camp listing shows the organization name, contact details, and current capacity. Users can easily select a camp to submit a donation offer. The map integration uses React Leaflet for smooth performance and accurate location tracking. This feature makes it convenient for donors to find nearby camps and participate in donations.""", ACCENT_BLUE)

# Slide 10: Heroes Leaderboard
add_image_slide(prs, "Heroes Leaderboard - Gamification",
    "screenshots/05_heroes_page.png",
    """The Heroes Leaderboard displays top donors ranked by their contributions and achievements. Each donor is shown with their rank, name, badge, points, blood group, and number of lives saved. The page features a badge system with seven levels: Starter, Bronze, Silver, Gold, Platinum, Diamond, and Legend. Donors earn badges based on the number of donations they complete. The leaderboard motivates donors by recognizing their contributions and showing their progress toward higher badges. Points are awarded for each approved donation, and the system calculates lives saved as donations multiplied by three. This gamification approach significantly increases donor engagement and retention.""", ACCENT_BLUE)

# Slide 11: User Dashboard
add_image_slide(prs, "User Dashboard - Personal Hub",
    "screenshots/06_user_dashboard.png",
    """The User Dashboard serves as the personal hub for donors, displaying their donation history and achievements. The dashboard shows an overview section with total donations, pending requests, approved requests, and rejected requests. Users can view their current badge, points, and progress toward the next badge level. The Quick Actions section provides buttons to submit blood donation offers or request emergency blood. The History section displays all past donations and requests with their status. Users can track their impact by seeing how many lives they've helped save. The dashboard is fully responsive and provides all necessary information at a glance.""", ACCENT_BLUE)

# Slide 12: Technology Stack - Frontend
add_paragraph_slide(prs, "Frontend Technology Stack",
"""LifeFlow's frontend is built with modern, industry-standard technologies. React 19 provides the foundation with its powerful component-based architecture and hooks for state management. Tailwind CSS 3 enables rapid UI development with utility-first styling, ensuring consistent design across all pages. Framer Motion adds smooth animations and transitions that enhance user experience. React Leaflet integrates interactive maps for camp location discovery. Zustand provides lightweight state management for handling application state efficiently. Axios serves as the HTTP client for seamless API communication with the backend. Together, these technologies create a fast, responsive, and engaging user interface that works flawlessly across all devices.""", ACCENT_PURPLE)

# Slide 13: Technology Stack - Backend
add_paragraph_slide(prs, "Backend Technology Stack",
"""The LifeFlow backend is powered by robust, scalable technologies. Node.js and Express provide a lightweight yet powerful server framework for building REST APIs. MySQL serves as the relational database for storing all application data securely. JWT (JSON Web Tokens) implements secure authentication, allowing users to maintain sessions without storing sensitive data on the server. Bcrypt handles password hashing with a 12-round salt, ensuring maximum security. Sequelize acts as the ORM (Object-Relational Mapping) layer, simplifying database operations and providing data validation. Geocoding APIs enable location-based services for camp mapping. This technology stack ensures scalability, security, and reliability for handling thousands of concurrent users.""", ACCENT_PURPLE)

# Slide 14: System Architecture
add_paragraph_slide(prs, "System Architecture",
"""LifeFlow follows a three-tier architecture pattern for optimal scalability and maintainability. The Presentation Layer consists of the React SPA (Single Page Application) with responsive design that works seamlessly on all devices. The Application Layer includes the Express REST API that handles all business logic, authentication, and data processing. The Data Layer uses MySQL with Sequelize ORM for efficient data persistence and retrieval. The system implements JWT-based authentication and authorization for secure access control. Role-based access control (RBAC) ensures that users can only access features appropriate for their role. Caching strategies optimize performance by reducing database queries. This architecture allows LifeFlow to scale horizontally and handle increasing user loads efficiently.""", ACCENT_PURPLE)

# Slide 15: Badge & Gamification System
add_paragraph_slide(prs, "Gamification - Badge System",
"""LifeFlow incorporates a comprehensive gamification system to motivate and reward donors. The badge system consists of seven progressive levels that donors can achieve based on their donation count. Starter Badge is awarded to new donors with zero donations. Bronze Badge is earned after the first donation. Silver Badge requires three donations. Gold Badge is achieved at five donations. Platinum Badge requires ten donations. Diamond Badge is earned at twenty donations. Legend Badge is the ultimate achievement, requiring fifty or more donations. Each badge comes with a unique emoji and color for visual recognition. Donors earn fifty points for each approved donation, and their points are displayed on the leaderboard. This system creates a sense of achievement and encourages regular participation.""", ACCENT_PURPLE)

# Slide 16: Admin Dashboard
add_paragraph_slide(prs, "Admin Dashboard - System Control",
"""The Admin Dashboard provides comprehensive system management capabilities for administrators. The overview section displays key metrics including total users, organizations, pending requests, and blood stock levels. Administrators can approve or reject blood donation requests and emergency blood requests. The user management section allows admins to view all registered users, their profiles, and donation history. The organization management section displays all registered organizations and their created camps. Administrators can review and approve donation camps before they become visible to donors. The analytics section provides insights into system usage, donation trends, and donor engagement. The badge and points system can be monitored to ensure proper gamification. This centralized control panel enables efficient system management and oversight.""", ACCENT_PURPLE)

# Slide 17: Organization Panel
add_paragraph_slide(prs, "Organization Panel - Camp Management",
"""The Organization Panel empowers blood banks and NGOs to manage donation camps effectively. Organizations can create new donation camps by providing camp name, description, location, date, time, and capacity. The camp creation form includes location selection on an interactive map for accurate positioning. Organizations can set the blood types needed and define the number of available slots. The camp management section displays all created camps with their status (pending, approved, or rejected). Organizations can view participant registrations for each camp and track attendance. Real-time notifications alert organizations when new donors register for their camps. Organizations can edit or cancel camps as needed. Analytics show camp statistics including total registrations, donation success rate, and donor feedback. This panel streamlines camp management and improves organizational efficiency.""", ACCENT_PURPLE)

# Slide 18: Security & Best Practices
add_paragraph_slide(prs, "Security & Best Practices",
"""LifeFlow implements industry-standard security practices to protect user data and system integrity. JWT token-based authentication ensures secure user sessions without storing sensitive data on the server. Bcrypt password hashing with 12-round salt provides maximum protection against brute-force attacks. Role-based access control (RBAC) ensures users can only access features appropriate for their role. Rate limiting on API endpoints prevents abuse and protects against DDoS attacks. CORS (Cross-Origin Resource Sharing) configuration restricts API access to authorized domains. Input validation and sanitization prevent SQL injection and XSS attacks. Error handling is implemented carefully to avoid exposing sensitive system information. All data transmission uses HTTPS encryption. Regular security audits and updates ensure the system remains protected against emerging threats.""", ACCENT_PURPLE)

# Slide 19: Benefits & Impact
add_paragraph_slide(prs, "Benefits & Real-World Impact",
"""LifeFlow delivers significant benefits to all stakeholders in the blood donation ecosystem. For patients and hospitals, the platform ensures faster access to blood during emergencies, potentially saving lives. For donors, the gamification system provides recognition and motivation to donate regularly. For organizations, the platform streamlines camp management and increases donor participation. The centralized system improves transparency and accountability in blood donation tracking. Data-driven insights help organizations optimize their donation campaigns. The community-building aspect creates a network of dedicated donors who feel valued and recognized. The platform reduces administrative burden on blood banks and hospitals. By increasing donation rates and improving efficiency, LifeFlow contributes to solving the global blood shortage crisis.""", ACCENT_PURPLE)

# Slide 20: Future Enhancements
add_paragraph_slide(prs, "Future Enhancements & Roadmap",
"""LifeFlow has an exciting roadmap of future enhancements to expand its capabilities. Mobile applications for iOS and Android will extend platform access to smartphone users. Push notifications will alert donors about urgent blood requests and nearby camps. AI-powered donor matching algorithms will optimize the matching process for better efficiency. Real-time chat support will provide immediate assistance to users. Advanced analytics and reporting will offer deeper insights into donation trends. Multi-language support will make the platform accessible to users worldwide. Integration with hospital management systems will streamline blood inventory management. Blockchain technology may be explored for enhanced transparency and security. Expansion to multiple regions and countries will scale the platform's impact. These enhancements will make LifeFlow the leading blood donation platform globally.""", ACCENT_PURPLE)

# Slide 21: Deployment & Scalability
add_paragraph_slide(prs, "Deployment & Scalability",
"""LifeFlow is designed for cloud deployment and horizontal scalability. The system can be deployed on major cloud platforms including AWS, Azure, or Google Cloud Platform. Auto-scaling capabilities automatically adjust server resources based on traffic demand. Content Delivery Networks (CDN) distribute static assets globally for faster loading. Database indexing and query optimization ensure efficient data retrieval even with large datasets. Redis caching reduces database load and improves response times. Load balancing distributes traffic across multiple servers for high availability. Monitoring and alerting systems track system health and notify administrators of issues. Docker containerization enables consistent deployment across different environments. CI/CD pipelines automate testing and deployment processes. This infrastructure ensures LifeFlow can handle millions of users while maintaining high performance and reliability.""", ACCENT_PURPLE)

# Slide 22: Development Team & Process
add_paragraph_slide(prs, "Development Team & Process",
"""LifeFlow was developed using modern development practices and methodologies. The frontend team used React, Tailwind CSS, and Framer Motion to create an engaging user interface. The backend team built a robust REST API using Node.js and Express with MySQL database. Development tools include Git for version control, Docker for containerization, and CI/CD pipelines for automation. The team follows Agile development methodology with regular sprints and iterations. Comprehensive testing including unit tests, integration tests, and end-to-end tests ensures code quality. Code reviews and pair programming maintain high code standards. Community feedback is actively incorporated into development decisions. Regular updates and maintenance keep the system secure and performant. The collaborative approach ensures LifeFlow meets user needs and maintains technical excellence.""", ACCENT_PURPLE)

# Slide 23: Conclusion
add_title_slide(prs, "Together We Save Lives", "Join LifeFlow - Make a Difference Today 🩸❤️")

# Save presentation
prs.save('LifeFlow_Final_Presentation.pptx')
print("✅ Final presentation created: LifeFlow_Final_Presentation.pptx")
print("📊 Total slides: 23")
print("🎨 Features: Screenshots, paragraph descriptions, colorful design")
print("📁 Location: C:\\xampp\\htdocs\\LifeFlow\\LifeFlow_Final_Presentation.pptx")
