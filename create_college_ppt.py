from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (Premium & Colorful)
PRIMARY_RED = RGBColor(220, 38, 38)
DARK_BG = RGBColor(15, 15, 25)
WHITE = RGBColor(255, 255, 255)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_GREEN = RGBColor(34, 197, 94)
ACCENT_PURPLE = RGBColor(139, 92, 246)
ACCENT_ORANGE = RGBColor(249, 115, 22)
TEXT_GRAY = RGBColor(75, 85, 99)

def add_title_slide(prs, title, subtitle, students=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Decorative shape
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_RED
    shape.line.color.rgb = PRIMARY_RED
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(68)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Student Details
    if students:
        student_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
        tf = student_box.text_frame
        for student in students:
            p = tf.add_paragraph()
            p.text = f"Presented By: {student['name']} ({student['enrollment']})"
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points, accent_color=PRIMARY_RED):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title Header
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = accent_color
    header.line.color.rgb = accent_color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(10)
    return slide

def add_image_info_slide(prs, title, image_path, description, use_case, tech_used, accent_color=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = accent_color
    header.line.color.rgb = accent_color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Image (Left)
    # Check if file exists, else use the raw media file from turn if it was copied
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(6.2), height=Inches(3.8))
    else:
        placeholder = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(6.2), Inches(3.8))
        placeholder.text = "Missing Image: " + os.path.basename(image_path)
    
    # Info (Right)
    info_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(2.8), Inches(4))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Feature Details:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = description
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "\nUse Case:"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = use_case
    p.font.size = Pt(14)
    
    # Languages/Tools (Bottom Full Width)
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(2))
    tf = tech_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Tech Stack Used:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    p = tf.add_paragraph()
    p.text = tech_used
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY
    
    return slide

# --- BEGIN SLIDES ---

# Student Details
team = [
    {"name": "Dhruv Solanki", "enrollment": "23001102031"},
    {"name": "Govind Desai", "enrollment": "23001102022"}
]

# Slide 1: Title
add_title_slide(prs, "LifeFlow", "Modern Blood Donation Management System", team)

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "LifeFlow is a full-stack platform designed to modernize blood donation processes.",
    "Goal: To create a bridge between willing donors and patients in critical need.",
    "Implements a secure, role-based ecosystem for Donors, Organizations, and Admins.",
    "Features real-time inventory tracking, emergency request systems, and gamified loyalty.",
    "Scalable architecture build using industry-standard modern web technologies.",
    "Focuses on accessibility, speed, and community impact."
], PRIMARY_RED)

# Slide 3: Problem Statement
add_content_slide(prs, "The Problem We Are Solving", [
    "Traditional blood donation systems lack real-time visibility during emergencies.",
    "Low retention: Donors often donate once but lack incentives to return.",
    "Logistical overhead for organizations managing manual camp registrations.",
    "Lack of verified status tracking for organizations and donor medical history.",
    "Patients struggle to find specific blood groups across different cities.",
    "Inefficient communication channels during critical hours."
], ACCENT_ORANGE)

# Slide 4: Tech Stack - Overview
add_content_slide(prs, "Our Technology Arsenal", [
    "Frontend Engine: React 19 for declarative and component-based UI.",
    "Design System: Tailwind CSS 3 for highly responsive and custom aesthetic.",
    "Animations: Framer Motion for premium user interactions and transitions.",
    "State Management: Zustand for efficient, high-performance global state control.",
    "Server-side: Node.js & Express - handling high-concurrency API requests.",
    "Data Persistence: MySQL with Sequelize ORM for structured relational data.",
    "Security Layers: JWT for session security and Bcrypt for data protection.",
    "Developer Tooling: Git, npm, Postman, and VS Code."
], ACCENT_PURPLE)

# Slide 5: Home Page - Hero Section
add_image_info_slide(prs, "The Landing Experience",
    "screenshots/01_home_hero.png",
    "A clean, professional landing page featuring a hero section with high-impact visuals.",
    "Designed to welcome users and provide clear navigation to essential features.",
    "React, Tailwind CSS, Framer Motion, Responsive Grid Layout.",
    ACCENT_BLUE)

# Slide 6: Home Page - Education
add_image_info_slide(prs, "Community Awareness",
    "screenshots/01_home_why.png",
    "Interactive section highlighting 'Why Donate Blood?' with curated medical facts.",
    "Used to educate new visitors and debunk myths about blood donation.",
    "SVG Icons, React Components, Accessible Typography.",
    ACCENT_BLUE)

# Slide 7: Login - Access Portal
add_image_info_slide(prs, "Secure Access Sanctuary",
    "screenshots/02_login.png",
    "A custom-designed login portal with secure authentication and role validation.",
    "Ensures that only authorized heroes and partners can access private data rooms.",
    "JWT (JSON Web Tokens), Express Middleware, Salted Password Hashing.",
    ACCENT_BLUE)

# Slide 8: Register - User Pathway
add_image_info_slide(prs, "The Onboarding Pipeline",
    "screenshots/03_register.png",
    "Dynamic registration page with a role-selection pathway (Donor vs Organization).",
    "Handles complex data entry for medical profiles vs institutional credentials.",
    "React Hook Form, Zod Schema Validation, Sequelize Model Logic.",
    ACCENT_BLUE)

# Slide 9: Discovery - Donation Camps
add_image_info_slide(prs, "Live Camp Navigator",
    "screenshots/04_camps.png",
    "A full-screen interactive mapping tool showing all active donation drives nearby.",
    "Used by donors to identify local events and join donation sessions instantly.",
    "React Leaflet, OpenStreetMap, Geolocation APIs.",
    ACCENT_BLUE)

# Slide 10: Gamification - Badge Tiers
add_image_info_slide(prs, "The Hero Progression",
    "screenshots/05_heroes_badges.png",
    "A loyalty-driven badge system ranging from Starter to Legend (0 to 2500+ points).",
    "Drives user engagement and repeat donations through social recognition.",
    "State-driven CSS styling, Milestone logic, Reward algorithms.",
    ACCENT_GREEN)

# Slide 11: Achievements - Personal Stats
add_image_info_slide(prs, "Individual Impact Profile",
    "screenshots/05_heroes_achievements.png",
    "Detailed profile view showing lives saved, donation count, and unlocked achievements.",
    "Provides donors with a sense of pride and a clear record of their health impact.",
    "Lucide React, Context-based data fetching, CSS Flexbox.",
    ACCENT_GREEN)

# Slide 12: Heroes - Real-time Analytics
add_image_info_slide(prs, "Global Impact Analytics",
    "screenshots/05_heroes_analytics.png",
    "Live statistics showing critical blood needs, total donors, and active requests.",
    "Provides transparency to the community and signals urgent needs in specific regions.",
    "Recharts, SQL View Aggregates, React useEffect hooks.",
    ACCENT_GREEN)

# Slide 13: Dashboard - Sanctuary Overview
add_image_info_slide(prs, "The Donor Command Center",
    "screenshots/06_dashboard_overview.png",
    "Personalized dashboard hub for donors to track their 'Impact Level' and actions.",
    "Centralized navigation for offering donations and making emergency requests.",
    "Zustand Global Store, Animated Progress Components.",
    ACCENT_PURPLE)

# Slide 14: Emergency - Rapid Response
add_image_info_slide(prs, "Emergency Request System",
    "screenshots/06_dashboard_emergency.png",
    "Fast-track form for patients and hospitals to broadcast urgent blood needs.",
    "Triggers system-wide alerts to match with donors of the required blood group.",
    "REST API POST, Validation Guards, MySQL Foreign Keys.",
    ACCENT_PURPLE)

# Slide 15: Participation - Giving Back
add_image_info_slide(prs, "Direct Donation Offering",
    "screenshots/06_dashboard_offer.png",
    "Interface allowing donors to volunteer for specific camps or emergency events.",
    "Allows organizations to prepare based on verified donor commitments.",
    "React Dropdowns, Relationship Mapping, Transactional SQL.",
    ACCENT_PURPLE)

# Slide 16: Journey - The Medical Timeline
add_image_info_slide(prs, "My Donation Journey",
    "screenshots/06_dashboard_journey.png",
    "A medical timeline showing every donation event and its verified completion status.",
    "Acts as a digital health record for donor participation over years.",
    "Timeline UI Pattern, CSS Grid, Array Mapping in React.",
    ACCENT_PURPLE)

# Slide 17: Admin - System Nexus
add_image_info_slide(prs, "The Nexus Control Panel",
    "screenshots/07_admin_monitor.png",
    "Holistic system monitor for administrators to oversee total supply vs demand.",
    "Used to manage critical shortages and monitor infrastructure health.",
    "Admin-only Route Guards, Sequelize Aggregate Queries.",
    PRIMARY_RED)

# Slide 18: Admin - Blood Request Validation
add_image_info_slide(prs, "Critical Need Authorization",
    "screenshots/07_admin_requests.png",
    "The validation portal where admins verify hospital requests before public broadcast.",
    "Prevents system abuse and ensures resources reach legitimate emergencies.",
    "Optimistic UI updates, Axios Put handlers.",
    PRIMARY_RED)

# Slide 19: Admin - Global Registry
add_image_info_slide(prs, "Universal Subject Registry",
    "screenshots/07_admin_users.png",
    "Complete directory of all registered donors with advanced search and filtering.",
    "Used by admins to identify and contact rare blood group donors during crises.",
    "SQL Pagination, Debounced Search, Restricted Data Views.",
    PRIMARY_RED)

# Slide 20: Admin - Partner Oversight
add_image_info_slide(prs, "Organizational Management",
    "screenshots/07_admin_orgs.png",
    "Verification pool for partner hospitals, blood banks, and NGO collaborations.",
    "Ensures all participating organizations meet safety and quality standards.",
    "Model Relationships, Status Flags, Audit Trails.",
    PRIMARY_RED)

# Slide 21: Admin - Donation Pipeline
add_image_info_slide(prs, "Pipeline Approval System",
    "screenshots/07_admin_offers.png",
    "The final step where donations are verified to award points and update levels.",
    "Secures the integrity of the gamification system by preventing fake entries.",
    "Atomic MySQL Updates, Point Calculation Logic.",
    PRIMARY_RED)

# Slide 22: Organization - The Hub
add_image_info_slide(prs, "Partner Command Center",
    "screenshots/08_org_dashboard.png",
    "Dedicated interface for NGOs to track their camp approvals and pending reviews.",
    "Empowers organizations to manage their own donation events independently.",
    "Modular React Components, Org-specific API scopes.",
    ACCENT_ORANGE)

# Slide 23: Organization - Camp Proposals
add_image_info_slide(prs, "Creating Donation Drives",
    "screenshots/08_org_create_camp.png",
    "Professional form for organizations to propose new camps for admin approval.",
    "Includes logistical details like slots, locations, and contact info.",
    "Complex Form Handling, React state validation.",
    ACCENT_ORANGE)

# Slide 24: Organization - Portfolio
add_image_info_slide(prs, "My Camp History",
    "screenshots/08_org_my_camps.png",
    "A clean dashboard list of all camps managed by the specific organization.",
    "Used to track donor turnout and status of ongoing donation drives.",
    "List Rendering, Conditional Badge Styling in Tailwind.",
    ACCENT_ORANGE)

# Slide 25: Closing - The Vision
add_title_slide(prs, "LifeFlow: The Future of Giving", "Connecting Hearts, Saving Lives.\nSubmission for College Project 2026", team)

# Save
target_file = 'LifeFlow_Final_Submission.pptx'
prs.save(target_file)
print(f"✅ Final College Presentation Created: {target_file}")
print("📊 Total Slides: 25 (All unique screenshots mapped)")
