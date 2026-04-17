from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_RED = RGBColor(220, 38, 38)
DARK_BG = RGBColor(20, 20, 30)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_GREEN = RGBColor(34, 197, 94)
ACCENT_PURPLE = RGBColor(139, 92, 246)
ACCENT_ORANGE = RGBColor(249, 115, 22)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add decorative shapes
    shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = PRIMARY_RED
    shape1.line.color.rgb = PRIMARY_RED
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points, bg_color=WHITE):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_RED
    title_shape.line.color.rgb = PRIMARY_RED
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_top = Inches(0.15)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BG
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = ACCENT_BLUE
    title_shape.line.color.rgb = ACCENT_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BG
        p.space_before = Pt(6)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "LifeFlow", "🩸 Blood Donation Management System 🩸")

# Slide 2: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "LifeFlow is a comprehensive digital platform designed to revolutionize blood donation management",
    "Connects donors, organizations, and patients in a seamless ecosystem",
    "Solves critical blood shortage problems through technology and gamification",
    "Provides real-time matching, tracking, and incentive systems",
    "Scalable solution for hospitals, blood banks, and NGOs",
    "Improves emergency response time and saves lives"
])

# Slide 3: Problem Statement
add_content_slide(prs, "The Problem We Solve", [
    "❌ Critical blood shortage in hospitals and medical centers worldwide",
    "❌ Difficulty finding willing donors during emergencies",
    "❌ No centralized platform for blood donation management",
    "❌ Lack of transparency in donation tracking and verification",
    "❌ Limited awareness about donation camps and opportunities",
    "❌ No incentive system to encourage regular donors"
])

# Slide 4: Our Solution
add_content_slide(prs, "Our Solution", [
    "✅ Centralized blood donation management platform",
    "✅ Easy donor registration with blood group and location tracking",
    "✅ Real-time blood request matching system",
    "✅ Location-based donation camp discovery with interactive maps",
    "✅ Gamification with badges, points, and leaderboard",
    "✅ Admin panel for request approval and system management"
])

# Slide 5: Key Features Overview
add_content_slide(prs, "Key Features", [
    "👤 Secure Authentication - Email/password with JWT tokens",
    "🗺️ Donation Camps - Interactive map with camp locations and details",
    "🩸 Blood Requests - Emergency blood request system with matching",
    "🎖️ Heroes Leaderboard - Gamified donor ranking and achievements",
    "📊 Admin Dashboard - Complete system management and analytics",
    "🏢 Organization Panel - Camp creation and event management"
])

# Slide 6: User Roles & Access
add_two_column_slide(prs, "User Roles & Permissions",
    "Donor Role", [
        "• Register and create profile",
        "• View available donation camps",
        "• Submit blood donation offers",
        "• Request emergency blood",
        "• Track donation history",
        "• Earn badges and points",
        "• View leaderboard ranking"
    ],
    "Admin & Organization", [
        "ADMIN:",
        "• Approve/reject requests",
        "• Manage all users",
        "• View system analytics",
        "",
        "ORGANIZATION:",
        "• Create donation camps",
        "• Manage camp schedules",
        "• Track registrations"
    ]
)

# Slide 7: Technology Stack - Frontend
add_content_slide(prs, "Frontend Technology Stack", [
    "⚛️ React 19 - Modern UI library with hooks and concurrent features",
    "🎨 Tailwind CSS 3 - Utility-first CSS framework for responsive design",
    "🎬 Framer Motion - Smooth animations and transitions",
    "🗺️ React Leaflet - Interactive maps for camp locations",
    "📦 Zustand - Lightweight state management solution",
    "🔗 Axios - HTTP client for API communication"
])

# Slide 8: Technology Stack - Backend
add_content_slide(prs, "Backend Technology Stack", [
    "🚀 Node.js & Express - Server framework for REST API",
    "🗄️ MySQL - Relational database for data persistence",
    "🔐 JWT (JSON Web Tokens) - Secure authentication mechanism",
    "🔒 Bcrypt - Password hashing with 12-round salt",
    "📊 Sequelize - ORM for database operations",
    "🗺️ Geocoding API - Location services for camp mapping"
])

# Slide 9: System Architecture
add_content_slide(prs, "System Architecture", [
    "🏗️ Three-Tier Architecture:",
    "   • Presentation Layer - React SPA with responsive design",
    "   • Application Layer - Express REST API with business logic",
    "   • Data Layer - MySQL database with Sequelize ORM",
    "",
    "🔄 Features:",
    "   • JWT-based authentication & authorization",
    "   • Role-based access control (RBAC)",
    "   • Caching for performance optimization"
])

# Slide 10: Login Page
add_content_slide(prs, "Login Page - Features & Design", [
    "🔐 Secure Authentication:",
    "   • Email and password validation",
    "   • JWT token generation on successful login",
    "   • Session persistence with localStorage",
    "",
    "🎨 Design Elements:",
    "   • Modern glassmorphism design",
    "   • Responsive mobile-friendly layout",
    "   • Error handling and validation messages",
    "   • Quick links to registration"
])

# Slide 11: Registration Page
add_content_slide(prs, "Registration Page - User Onboarding", [
    "👤 Information Collection:",
    "   • Full name, email, password",
    "   • Blood group selection",
    "   • Age and location details",
    "   • Role selection (Donor/Organization)",
    "",
    "🏢 Organization-Specific Fields:",
    "   • Organization name",
    "   • Contact phone number",
    "   • Organization address",
    "   • Verification process"
])

# Slide 12: Home Page
add_content_slide(prs, "Home Page - Welcome & Engagement", [
    "🎯 Hero Section:",
    "   • Compelling call-to-action",
    "   • Impact statistics and metrics",
    "   • Quick navigation buttons",
    "",
    "📊 Content Sections:",
    "   • Featured heroes and top donors",
    "   • How the system works explanation",
    "   • Testimonials and success stories",
    "   • Responsive design for all devices"
])

# Slide 13: User Dashboard
add_content_slide(prs, "User Dashboard - Personal Hub", [
    "📊 Overview Section:",
    "   • Total donations and requests count",
    "   • Pending and approved requests",
    "   • Current badge and points display",
    "",
    "🎯 Quick Actions:",
    "   • Submit blood donation offer",
    "   • Request emergency blood",
    "   • View donation history",
    "   • Track achievements and progress"
])

# Slide 14: Donation Camps Page
add_content_slide(prs, "Donation Camps - Location Discovery", [
    "🗺️ Interactive Map Features:",
    "   • Real-time camp locations",
    "   • Zoom and pan functionality",
    "   • Location-based filtering",
    "",
    "📋 Camp Information:",
    "   • Camp name and description",
    "   • Date, time, and duration",
    "   • Address and contact details",
    "   • Available slots and blood types needed"
])

# Slide 15: Heroes Leaderboard
add_content_slide(prs, "Heroes Leaderboard - Gamification", [
    "🏆 Ranking System:",
    "   • Top donors displayed with rank",
    "   • Points-based scoring system",
    "   • Lives saved calculation (donations × 3)",
    "",
    "🎖️ Badge Progression:",
    "   • Starter (0) → Bronze (1) → Silver (3)",
    "   • Gold (5) → Platinum (10) → Diamond (20)",
    "   • Legend (50+) - Ultimate achievement",
    "   • Motivates regular donations"
])

# Slide 16: Admin Dashboard
add_content_slide(prs, "Admin Dashboard - System Control", [
    "📊 Overview & Analytics:",
    "   • Total users and organizations count",
    "   • Pending requests and donations",
    "   • Blood stock levels by type",
    "   • System statistics and metrics",
    "",
    "✅ Management Functions:",
    "   • Approve/reject donation requests",
    "   • Manage users and organizations",
    "   • Review blood requests",
    "   • Monitor badge and points system"
])

# Slide 17: Organization Panel
add_content_slide(prs, "Organization Panel - Camp Management", [
    "🏢 Camp Creation:",
    "   • Create new donation camps",
    "   • Set date, time, and location",
    "   • Define capacity and blood types needed",
    "   • Add detailed description",
    "",
    "📊 Management & Tracking:",
    "   • View participant registrations",
    "   • Track camp statistics",
    "   • Receive notifications for new registrations",
    "   • Edit or cancel camps"
])

# Slide 18: Badge & Gamification System
add_content_slide(prs, "Gamification - Badge System", [
    "🌟 Starter Badge - 0 donations",
    "🥉 Bronze Badge - 1 donation",
    "🥈 Silver Badge - 3 donations",
    "🥇 Gold Badge - 5 donations",
    "💎 Platinum Badge - 10 donations",
    "👑 Legend Badge - 50+ donations",
    "",
    "Points System: 50 points per approved donation"
])

# Slide 19: Security & Best Practices
add_content_slide(prs, "Security & Best Practices", [
    "🔐 Authentication & Authorization:",
    "   • JWT token-based authentication",
    "   • Bcrypt password hashing (12 rounds)",
    "   • Role-based access control (RBAC)",
    "",
    "🛡️ API Security:",
    "   • Rate limiting on endpoints",
    "   • CORS configuration",
    "   • Input validation and sanitization",
    "   • Error handling without exposing sensitive data"
])

# Slide 20: Benefits & Impact
add_content_slide(prs, "Benefits & Real-World Impact", [
    "💚 Life-Saving Impact:",
    "   • Efficient blood management reduces shortages",
    "   • Faster emergency response times",
    "   • Increased donation rates through gamification",
    "",
    "🤝 Community Building:",
    "   • Creates community of regular donors",
    "   • Recognizes and rewards contributions",
    "   • Builds trust between donors and organizations",
    "   • Scalable to multiple regions and countries"
])

# Slide 21: Future Enhancements
add_content_slide(prs, "Future Enhancements & Roadmap", [
    "📱 Mobile Applications:",
    "   • Native iOS and Android apps",
    "   • Push notifications for urgent requests",
    "   • Offline functionality",
    "",
    "🤖 Advanced Features:",
    "   • AI-powered donor matching algorithm",
    "   • Real-time chat support system",
    "   • Advanced analytics and reporting",
    "   • Multi-language support"
])

# Slide 22: Deployment & Scalability
add_content_slide(prs, "Deployment & Scalability", [
    "☁️ Cloud Infrastructure:",
    "   • Deployed on cloud platforms (AWS/Azure/GCP)",
    "   • Auto-scaling for traffic spikes",
    "   • CDN for static assets",
    "",
    "📈 Performance Optimization:",
    "   • Database indexing and query optimization",
    "   • Caching strategies (Redis)",
    "   • Load balancing across servers",
    "   • Monitoring and alerting systems"
])

# Slide 23: Team & Development
add_content_slide(prs, "Development Team & Process", [
    "👨‍💻 Tech Stack Summary:",
    "   • Frontend: React, Tailwind CSS, Framer Motion",
    "   • Backend: Node.js, Express, MySQL",
    "   • Tools: Git, Docker, CI/CD pipelines",
    "",
    "🔄 Development Methodology:",
    "   • Agile development approach",
    "   • Regular testing and quality assurance",
    "   • Continuous integration and deployment",
    "   • Community feedback integration"
])

# Slide 24: Conclusion & Call to Action
add_title_slide(prs, "Together We Save Lives", "Join LifeFlow - Make a Difference Today 🩸❤️")

# Save presentation
prs.save('LifeFlow_Detailed_Presentation.pptx')
print("✅ Detailed presentation created: LifeFlow_Detailed_Presentation.pptx")
print("📊 Total slides: 24")
print("🎨 Features: Colorful design, detailed descriptions, comprehensive coverage")
