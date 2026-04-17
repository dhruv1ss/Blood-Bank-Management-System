from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_RED = RGBColor(220, 38, 38)  # Red
DARK_BG = RGBColor(20, 20, 30)  # Dark blue-black
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_GREEN = RGBColor(34, 197, 94)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_RED
    title_shape.line.color.rgb = PRIMARY_RED
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_top = Inches(0.15)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BG
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(10)
    
    return slide

def add_image_slide(prs, title, image_path, description):
    """Add a slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = ACCENT_BLUE
    title_shape.line.color.rgb = ACCENT_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.5)
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = description
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BG
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "LifeFlow", "Blood Donation Management System")

# Slide 2: Overview
add_content_slide(prs, "What is LifeFlow?", [
    "🩸 A comprehensive blood donation management platform",
    "🏥 Connects donors, organizations, and patients",
    "📍 Location-based donation camp discovery",
    "🎖️ Gamification system with badges and points",
    "⚡ Real-time request matching and notifications",
    "🔐 Secure authentication and role-based access"
])

# Slide 3: Problem Statement
add_content_slide(prs, "The Problem", [
    "❌ Blood shortage in hospitals and medical centers",
    "❌ Difficult to find willing donors in emergencies",
    "❌ No centralized platform for blood donation management",
    "❌ Lack of transparency in donation tracking",
    "❌ Limited awareness about donation camps",
    "❌ No incentive system for regular donors"
])

# Slide 4: Solution
add_content_slide(prs, "Our Solution", [
    "✅ Centralized blood donation management system",
    "✅ Easy donor registration and profile management",
    "✅ Real-time blood request matching",
    "✅ Donation camp discovery with maps",
    "✅ Gamification with badges and leaderboard",
    "✅ Admin panel for request approval and management"
])

# Slide 5: Key Features
add_content_slide(prs, "Key Features", [
    "👤 User Authentication - Secure login/registration",
    "🗺️ Donation Camps - Location-based camp discovery",
    "🩸 Blood Requests - Emergency blood request system",
    "🎖️ Heroes Leaderboard - Gamified donor ranking",
    "📊 Admin Dashboard - Request and donation management",
    "🏢 Organization Panel - Camp creation and management"
])

# Slide 6: User Roles
add_content_slide(prs, "User Roles", [
    "👤 DONOR - Register, donate blood, track donations, earn badges",
    "🏢 ORGANIZATION - Create donation camps, manage events",
    "👨‍💼 ADMIN - Approve requests, manage users, oversee system",
    "🔐 Role-based access control for security",
    "📱 Different dashboards for each role",
    "🔔 Role-specific notifications and alerts"
])

# Slide 7: Technology Stack - Frontend
add_content_slide(prs, "Technology Stack - Frontend", [
    "⚛️ React 19 - Modern UI library",
    "🎨 Tailwind CSS 3 - Utility-first styling",
    "🎬 Framer Motion - Smooth animations",
    "🗺️ React Leaflet - Interactive maps",
    "📦 Zustand - State management",
    "🔗 Axios - HTTP client for API calls"
])

# Slide 8: Technology Stack - Backend
add_content_slide(prs, "Technology Stack - Backend", [
    "🚀 Node.js & Express - Server framework",
    "🗄️ MySQL - Relational database",
    "🔐 JWT - Secure authentication",
    "🔒 Bcrypt - Password hashing",
    "📊 Sequelize - ORM for database",
    "🗺️ Geocoding API - Location services"
])

# Slide 9: Architecture Overview
add_content_slide(prs, "System Architecture", [
    "🏗️ Three-tier architecture",
    "📱 Frontend - React SPA with responsive design",
    "🔌 Backend - Express REST API",
    "🗄️ Database - MySQL with Sequelize ORM",
    "🔄 Real-time updates with caching",
    "🔐 JWT-based authentication & authorization"
])

# Slide 10: Login Page
add_content_slide(prs, "Login Page Features", [
    "🔐 Secure email and password authentication",
    "🎨 Modern glassmorphism design",
    "📱 Responsive mobile-friendly layout",
    "🔄 Remember me functionality",
    "🔗 Quick links to registration",
    "⚠️ Error handling and validation"
])

# Slide 11: Registration Page
add_content_slide(prs, "Registration Page Features", [
    "👤 User information collection",
    "🩸 Blood group selection",
    "📍 Location and address details",
    "🏢 Organization-specific fields",
    "✅ Form validation and error messages",
    "🔐 Secure password requirements"
])

# Slide 12: Home Page
add_content_slide(prs, "Home Page Features", [
    "🎯 Hero section with call-to-action",
    "📊 Statistics and impact metrics",
    "🎖️ Featured heroes and top donors",
    "🏥 How it works explanation",
    "📱 Responsive design for all devices",
    "🔗 Quick navigation to key features"
])

# Slide 13: User Dashboard
add_content_slide(prs, "User Dashboard Features", [
    "📊 Overview of donations and requests",
    "🎖️ Badge and points display",
    "🩸 Submit blood donation offers",
    "🏥 Request emergency blood",
    "📜 History of all transactions",
    "🔔 Notifications and alerts"
])

# Slide 14: Donation Camps Page
add_content_slide(prs, "Donation Camps Features", [
    "🗺️ Interactive map with camp locations",
    "📍 Filter camps by location and date",
    "📋 Detailed camp information",
    "🎯 Easy camp selection for donations",
    "⏰ Schedule and timing details",
    "📞 Contact information for camps"
])

# Slide 15: Heroes Leaderboard
add_content_slide(prs, "Heroes Leaderboard Features", [
    "🏆 Top donors ranking system",
    "🎖️ Badge progression (Starter → Legend)",
    "📊 Points and achievements display",
    "🩸 Lives saved calculation",
    "🌍 Location-based statistics",
    "🎯 Gamification to encourage donations"
])

# Slide 16: Admin Dashboard
add_content_slide(prs, "Admin Dashboard Features", [
    "📊 System overview and statistics",
    "✅ Approve/reject donation requests",
    "👥 Manage users and organizations",
    "🏥 Review blood requests",
    "🎖️ Monitor badge and points system",
    "📈 Analytics and reporting"
])

# Slide 17: Organization Panel
add_content_slide(prs, "Organization Panel Features", [
    "🏢 Create and manage donation camps",
    "📅 Schedule camp dates and times",
    "📍 Set camp locations with maps",
    "👥 Track participant registrations",
    "📊 View camp statistics",
    "🔔 Receive notifications for new registrations"
])

# Slide 18: Badge System
add_content_slide(prs, "Gamification - Badge System", [
    "🌟 Starter - 0 donations",
    "🥉 Bronze - 1 donation",
    "🥈 Silver - 3 donations",
    "🥇 Gold - 5 donations",
    "💎 Platinum - 10 donations",
    "👑 Legend - 50+ donations"
])

# Slide 19: Security Features
add_content_slide(prs, "Security & Best Practices", [
    "🔐 JWT token-based authentication",
    "🔒 Bcrypt password hashing (12 rounds)",
    "🛡️ Role-based access control (RBAC)",
    "⚡ Rate limiting on API endpoints",
    "🔄 CORS configuration for security",
    "📝 Input validation and sanitization"
])

# Slide 20: Benefits & Impact
add_content_slide(prs, "Benefits & Impact", [
    "💚 Saves lives through efficient blood management",
    "🤝 Builds community of regular donors",
    "📈 Increases donation rates through gamification",
    "⏱️ Reduces emergency response time",
    "🌍 Scalable to multiple regions",
    "📊 Data-driven insights for organizations"
])

# Slide 21: Future Enhancements
add_content_slide(prs, "Future Enhancements", [
    "📱 Mobile app for iOS and Android",
    "🤖 AI-powered donor matching algorithm",
    "💬 Real-time chat support system",
    "📊 Advanced analytics and reporting",
    "🌐 Multi-language support",
    "🔔 Push notifications for urgent requests"
])

# Slide 22: Conclusion
add_title_slide(prs, "Thank You!", "Together, we save lives through LifeFlow 🩸❤️")

# Save presentation
prs.save('LifeFlow_Presentation.pptx')
print("✅ Presentation created successfully: LifeFlow_Presentation.pptx")
